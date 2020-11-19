""" Allows for the redaction of words or phrases from a variety of file formats """
__author__ = "Evan Scott"
__email__  = "escott1367@gmail.com"
__status__ = "Prototype"

from os import write
import fitz 
import re
import copy
from PIL import Image, ImageDraw
from io import BytesIO
import requests
from django.conf import settings
import time

from requests.api import head 

from azure.cognitiveservices.vision.computervision import ComputerVisionClient
from azure.cognitiveservices.vision.computervision.models import OperationStatusCodes
from azure.cognitiveservices.vision.computervision.models import VisualFeatureTypes
from msrest.authentication import CognitiveServicesCredentials
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.shape import WD_INLINE_SHAPE
from docx.oxml.shared import OxmlElement
from docx.text.run import Run
import docx.dml.color
from pptx import Presentation
from pptx.dml.color import RGBColor, ColorFormat
from defusedxml.ElementTree import parse
from openpyxl import load_workbook
from openpyxl.styles import PatternFill, Font
import csv

M_VISION_KEY = getattr(settings, "M_VISION_KEY")
M_VISION_ENDPOINT = getattr(settings,"M_VISION_ENDPOINT")

class Redactor: 
    
    # constructor 
    def __init__(self, path): 
        self.path = path

    # uses the xml for the highlighting
    def set_hightlight_xml(self, run):
        rpr = run._r.get_or_add_rPr()
        highlight = OxmlElement("a:highlight")
        srgbClr = OxmlElement("a:srgbClr")
        setattr(srgbClr, "val", WD_COLOR_INDEX.BLACK)
        highlight.append(srgbClr)
        rpr.append(highlight)
        return run         

    def add_run_styles(self, new_run, old_run, pptx):
        
        new_run.font.name = old_run.font.name
        new_run.font.size = old_run.font.size
        
        if not pptx:
            color = new_run.font.color
            color.rgb = old_run.font.color.rgb
            new_run.bold = old_run.bold
            new_run.italic = old_run.italic
            new_run.style = old_run.style
            new_run.underline = old_run.underline
            new_run.font.all_caps = old_run.font.all_caps
            new_run.font.complex_script = old_run.font.complex_script
            new_run.font.double_strike = old_run.font.double_strike
            new_run.font.emboss = old_run.font.emboss
            new_run.font.highlight_color = old_run.font.highlight_color
            new_run.font.imprint = old_run.font.imprint 
            new_run.font.math = old_run.font.math
            new_run.font.outline = old_run.font.outline
            new_run.font.shadow = old_run.font.shadow
            new_run.font.small_caps = old_run.font.small_caps
            new_run.font.snap_to_grid = old_run.font.snap_to_grid
            new_run.font.strike = old_run.font.strike
            new_run.font.subscript = old_run.font.subscript
            new_run.font.superscript = old_run.font.superscript
        else:
            new_run.font.bold = old_run.font.bold
            color = old_run.font.color
            new_run.font.color._color = color._color
            new_run.font.color._xFill = color._xFill
            fill = old_run.font.fill
            new_run.font.fill._xPr = fill._xPr
            new_run.font.fill._fill = fill._fill
            new_run.font.italic = old_run.font.italic
            new_run.font.underline = old_run.font.underline

        return new_run


    def paragraph_rewrite(self, paragraph, phrase):
        curr_runs = copy.copy(paragraph.runs)    # copy the paragraph and clear it 
        new_paragraph = paragraph.clear()
        for i in range(len(curr_runs)): # Find which run contains the phrase
            if phrase.casefold() in curr_runs[i].text.casefold(): # The phrase to redact is in this run 
                # Replace the word with the dash_word to redact the phrase
                dash_word = "#"*len(phrase)
                text = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, curr_runs[i].text, flags=re.IGNORECASE)
                curr_runs[i].text = text
                # Split by the dash_word to add highlighting 
                words = re.split('(#+)', curr_runs[i].text)
                text_string = ""
                #Search for the redacted word in the run 
                for word in words:
                    # if the word is the redacted word create new run and black it out 
                    if word == dash_word:
                        # if the new_run has words attached to it add it to the paragraph and create a new 
                        # run for the redacted word
                        if text_string != "":
                            new_run = new_paragraph.add_run(text_string)
                            new_run = self.add_run_styles(new_run, curr_runs[i], False)
                        # If it is a fresh_run just create a new run containing only the redacted word and add it
                        text_string = ""
                        new_run = new_paragraph.add_run(dash_word)
                        new_run = self.add_run_styles(new_run, curr_runs[i], False)
                        new_run.font.highlight_color = WD_COLOR_INDEX.BLACK
                    # else just add the word to the existing text
                    else:
                        text_string += word
                # Deals with the remainder of the run after the phrase has been found
                if text_string != "":
                    new_run = new_paragraph.add_run(text_string)
                    new_run = self.add_run_styles(new_run, curr_runs[i], False)
            else:
                #Append the run as there is nothing to change
                #TODO cannot recognize the highlight color of the doc might need to go into lxml
                new_run = new_paragraph.add_run(curr_runs[i].text)
                new_run = self.add_run_styles(new_run, curr_runs[i], False)
        return new_paragraph

    def powerpoint_rewrite(self, paragraph, phrase):
        curr_runs = copy.copy(paragraph.runs)
        new_paragraph = paragraph.clear()
        for i, run in enumerate(curr_runs):
            # it is a tuple 
            if phrase.casefold() in run.text.casefold(): # The phrase to redact is in this run 
                # Replace the word with the dash_word to redact the phrase
                dash_word = "#"*len(phrase)
                text = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, run.text, flags=re.IGNORECASE)
                run.text = text
                # Split by the dash_word to add highlighting 
                words = re.split('(#+)', run.text)
                text_string = ""
                #Search for the redacted word in the run 
                for word in words:
                    if word == dash_word:
                        # if the new_run has words attached to it add it to the paragraph and create a new 
                        # run for the redacted word
                        if text_string != "":
                            new_run = new_paragraph.add_run()
                            new_run.text = text_string
                            new_run = self.add_run_styles(new_run, curr_runs[i], True)
                        # If it is a fresh_run just create a new run containing only the redacted word and add it
                        text_string = ""
                        new_run = new_paragraph.add_run()
                        new_run.text = dash_word
                        new_run = self.add_run_styles(new_run, curr_runs[i], True)
                        
                        color_format = ColorFormat(RGBColor(0,0,0),RGBColor(0,0,0))
                        print(str(color_format.__dict__))
                        new_run.font.fill.rgb = RGBColor(0,0,0)
                        new_run.font.fill._xPr = RGBColor(0,0,0)
                        new_run.font.fill._fill = RGBColor(0,0,0)
                        new_run.font.fill._fore_color = color_format
                        new_run.font.fill._back_color = color_format
                    # else just add the word to the existing text
                    else:
                        text_string += word
                if text_string != "":
                    new_run = new_paragraph.add_run()
                    new_run.text = text_string
                    new_run = self.add_run_styles(new_run, curr_runs[i], True)
            else:
                #Append the run as there is nothing to change
                #TODO cannot recognize the highlight color of the doc might need to go into lxml
                new_run = new_paragraph.add_run()
                new_run.text = run.text
                new_run = self.add_run_styles(new_run, curr_runs[i], True)
        return new_paragraph

    def img_overwrite(self, im, redacted_lines):
        # img
        ocr_url = M_VISION_ENDPOINT + "vision/v3.1/read/analyze"
        img_url = self.path
        headers = {'Ocp-Apim-Subscription-Key': M_VISION_KEY }
        data = {'url' :img_url}

        response = requests.post(ocr_url, headers=headers, json = data)
        response.raise_for_status() # checks for invalid request 
        # Need 2 API calls one for processing, one for retrieving 
        operation_url = response.headers["Operation-Location"]
        analysis = {}
        poll = True
        while(poll):
            response_final = requests.get(response.headers["Operation-Location"], headers=headers)
            analysis = response_final.json()
            time.sleep(1)
            if ("analyzeResult" in analysis):
                poll = False
            if("status" in analysis and analysis['status'] == 'failed'):
                poll = False

        if analysis['status'] != 'failed':
            for text_result in analysis['analyzeResult']['readResults']:
                for line in text_result['lines']:
                    for word in line['words']:
                        for phrase in redacted_lines:
                            if phrase.casefold() == word['text'].casefold():
                                loc = word['boundingBox']
                                # 
                                draw = ImageDraw.Draw(im)
                                # [x0,y0,x1,y1] or [(x0,y0), (x1,y1)] format
                                #draw.rectangle([(loc[0], loc[1]), (loc[4], loc[5])], fill=(0,0,0))
                                draw.polygon([(loc[0],loc[1]),(loc[2],loc[3]),(loc[4],loc[5]),(loc[6],loc[7])] ,fill=(0,0,0))
        return im

    def local_img_overwrite(self, im, redacted_lines):
        bin_img = BytesIO()
        im.save(bin_img, format='PNG')
        #im.close()
        image_data = bin_img.getvalue()
        bin_img.close()
        ocr_url = M_VISION_ENDPOINT + "vision/v3.1/read/analyze"
        headers = {'Ocp-Apim-Subscription-Key': M_VISION_KEY, 'Content-Type': 'application/octet-stream'}
        response = requests.post(ocr_url, headers=headers, data=image_data)
        response.raise_for_status() # checks for invalid request 
        analysis = {}
        poll = True
        while(poll):
            response_final = requests.get(response.headers["Operation-Location"], headers=headers)
            analysis = response_final.json()
            time.sleep(1)
            if ("analyzeResult" in analysis):
                poll = False
            if("status" in analysis and analysis['status'] == 'failed'):
                poll = False

        if analysis['status'] != 'failed':
            for text_result in analysis['analyzeResult']['readResults']:
                for line in text_result['lines']:
                    for word in line['words']:
                        for phrase in redacted_lines:
                            if phrase.casefold() == word['text'].casefold():
                                loc = word['boundingBox']
                                # 
                                draw = ImageDraw.Draw(im)
                                # [x0,y0,x1,y1] or [(x0,y0), (x1,y1)] format
                                #draw.rectangle([(loc[0], loc[1]), (loc[4], loc[5])], fill=(0,0,0))
                                draw.polygon([(loc[0],loc[1]),(loc[2],loc[3]),(loc[4],loc[5]),(loc[6],loc[7])] ,fill=(0,0,0))
        return im

    # Redacts document based on file extension 
    def redaction(self): 
        redacted_lines = ["powerpoint","ipsum", "phrases to redact", "over", "test phrase2", "jumps", "Yukon", "lazy", "computer", "canada", "website"]
        extension = self.path.split(".")[-1]
        
        # For docx
        if extension == "docx":
            doc = Document(self.path)
            # Clear one paragraph at a time 
            for paragraph in doc.paragraphs:
                for phrase in redacted_lines:
                    if phrase.casefold() in paragraph.text.casefold(): # There is something to redact
                       paragraph = self.paragraph_rewrite(paragraph, phrase)
        
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            for phrase in redacted_lines:
                                if phrase.casefold() in paragraph.text.casefold():
                                    paragraph = self.paragraph_rewrite(paragraph, phrase)
            # redacts images from docx files if needed
            for shape in doc.inline_shapes:
                if shape.type == WD_INLINE_SHAPE.PICTURE:
                    blip = shape._inline.graphic.graphicData.pic.blipFill.blip # gets <a:blip> element
                    r_id = blip.embed
                    document_part = doc.part
                    image_part = document_part.related_parts[r_id]
                    ext_image = Image.open(BytesIO(image_part))
                    if min(ext_image.size) > 50:
                        new_img = self.local_img_overwrite(ext_image, redacted_lines)
                        bin_img = BytesIO()
                        new_img.save(bin_img, format='PNG')
                        image_data = bin_img.getvalue()
                        bin_img.close()
                        image_part._blob = image_data
            
            doc.save('./redacted.docx')
            return('redacted.docx')

        # For pdfs
        elif extension == "pdf":
            # opening the pdf 
            doc = fitz.open(self.path) 
            # iterating through pages 
            for page in doc: 
                # _wrapContents is needed for fixing alignment issues with rect boxes 
                page._wrapContents() 
                                    
                for phrase in redacted_lines: 
                    areas = page.searchFor(phrase) 
                    for area in areas:
                        anot = page.addRedactAnnot(area, fill = (0,0,0))
                        r = anot.rect
                        r.y1 = r.y0 + r.height * .9
                        r.y0 = r.y1 - r.height * .9
                        anot.setRect(r)
                        anot.update()

                # applying the redaction 
                page.apply_redactions()
                # Now cycle through the pages images (list of lists)
                # The same image could be referenced multiple times so will want to filter this
                for image in page.getImageList(full=True):
                    # First turn the old image into a redacted version of itself
                    xref = image[0] # image contains [0] = xref, [1]=smask, [2]=width, [3]=height,...
                    ext_img = doc.extractImage(xref)
                    ext_image = Image.open(BytesIO(ext_img['image']))
                    if min(ext_image.size) > 50:    # Image needs to be at least this big to run through ocr
                        new_img = self.local_img_overwrite(ext_image, redacted_lines) # corrects returns redacted version of image
                        # Now set remove the old image and insert the new one
                        img_rect = page.getImageBbox(image[7])
                        page.addRedactAnnot(img_rect)
                        page.apply_redactions()
                        # Get new image into an insertable format
                        bin_img = BytesIO()
                        new_img.save(bin_img, format='PNG')
                        image_data = bin_img.getvalue()
                        bin_img.close()
                        page.insertImage(img_rect, stream=image_data)  

            # saving it to a new pdf 
            doc.save('redacted22.pdf', garbage=3, deflate = True)
            return('redacted22.pdf')
        
        # For images you need to feed it single words not phrases or it will not work
        # invert the image and dilate it, merges the letters a bit then use contours to find each word separately
        elif extension == "jpg" or extension == "png" or extension == "jpeg":
            response = requests.get(self.path)
            im = Image.open(BytesIO(response.content))
            im = self.img_overwrite(im, redacted_lines)
            save_loc = 'redactedpoly.' + extension
            im.save(save_loc)
            return(save_loc)

        # For powerpoints 
        elif extension == "pptx":
            presentation =  Presentation(self.path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            for phrase in redacted_lines:
                                if paragraph is not None  and phrase.casefold() in paragraph.text.casefold():
                                    paragraph = self.powerpoint_rewrite(paragraph, phrase)

            presentation.save('redacted-powerpoint.pptx')
            return('redacted-powerpoint.pptx')

        # For plain text files
        elif extension == "txt":
            with open(self.path, "r") as txt_file:
                new_file = open('redacted.txt', mode ='w')
                for line in txt_file.readlines():
                    for phrase in redacted_lines:
                        if phrase in line:
                            dash_word = "#"*len(phrase)
                            new_text = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, line, flags=re.IGNORECASE)
                            new_file.write(new_text)
                new_file.close()
                return('redacted.txt')

        # for csv files 
        elif extension == "csv":
            lines = list()
            with open(self.path) as csvfile:
                reader = csv.reader(csvfile, delimiter= ' ', quotechar = '|')
                for row in reader:
                    for field in row:
                        for phrase in redacted_lines:
                            if phrase in field: # Get rid of the phrase if it exists before adding the row
                                dash_word = "#"*len(phrase)
                                field = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, field, flags=re.IGNORECASE)
                            lines.append(row)
            with open('redacted-csv.csv', 'w') as writeFile:
                writer = csv.writer(writeFile)
                writer.writerow(lines)
            return ('redacted-csv.csv')

        # For xls files
        elif extension == "xlsx":
            workbook = load_workbook(self.path)
            worksheets = workbook.worksheets
            for sheet in worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        for phrase in redacted_lines:
                            if phrase in str(cell.value):
                                # Currently blacks out entire cell may not want to do this or at least set a flag
                                dash_word = "#"*len(phrase)
                                cell.value = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, str(cell.value), flags=re.IGNORECASE)
                                cell.fill = PatternFill(bgColor="000000", fill_type = "solid")
                                cell.font = Font(color = '000000')
            workbook.save('redacted.xlsx')
            return('redacted.xlsx')
                        

        # For xml files - need to use defusedxml to avoid a bomb
        elif extension == 'xml':
            tree = parse(self.path)
            for elem in tree.iter():
                for phrase in redacted_lines:
                    if elem.text is not None and phrase in elem.text:
                        dash_word = "#"*len(phrase)
                        elem.text = re.sub(r'\b%s\b' % re.escape(phrase), dash_word, elem.text, flags=re.IGNORECASE)
            tree.write('redacted.xml')
            return('redacted.xml')